from __future__ import annotations

import csv
import json
from datetime import date, timedelta
from pathlib import Path
from typing import Any

from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .metrics import add_days, latest_progress, project_progress, task_end_date, normalize_date, parse_int
from .models import ArchiveItem, DailyLog, InboxTask, ProgressEntry, Project, Task, Workspace, new_id, to_dict, today, APP_VERSION

LOCAL_KEYS = ("project-desk-local-v5", "project-desk-local-v4", "project-desk-v3")
RISKS = {"H", "M", "L"}
STATUSES = {"Open", "Ongoing", "Closed"}
ARCHIVE_TYPES = {"实验数据", "汇报PPT", "会议纪要", "图片截图", "交付版本", "其他"}
INBOX_STATUSES = {"待归档", "待处理", "已转项目任务", "已转任务", "已归档到项目", "已归档", "已新建项目", "已建项目", "已忽略"}


def pick(source: dict[str, Any], *keys: str, default: Any = None) -> Any:
    for key in keys:
        if key in source and source[key] is not None:
            return source[key]
    return default


def clamp_progress(value: Any) -> int:
    try:
        number = round(float(value))
    except (TypeError, ValueError):
        return 0
    return max(0, min(100, number))


def clean_date(value: Any, fallback: str | None = None) -> str:
    normalized = normalize_date(str(value or ""))
    if normalized:
        return normalized
    return fallback if fallback is not None else today()


def unwrap_payload(raw: Any) -> tuple[dict[str, Any], str]:
    if not isinstance(raw, dict):
        raise ValueError("JSON 顶层必须是对象。")
    for key in LOCAL_KEYS:
        if key in raw:
            value = raw[key]
            return (json.loads(value) if isinstance(value, str) else value), f"localStorage {key}"
    if isinstance(raw.get("payload"), dict):
        return raw["payload"], "Supabase legacy payload"
    if isinstance(raw.get("workspace"), dict):
        return raw["workspace"], "wrapped workspace"
    if isinstance(raw.get("data"), dict):
        return raw["data"], "wrapped data"
    if isinstance(raw.get("projects"), list):
        return raw, "workspace"
    if "tasks" in raw or "dailyLogs" in raw or "daily_logs" in raw:
        return {"projects": [raw]}, "single project"
    raise ValueError("没有识别到 workspace、projects、payload、data 或单项目结构。")


def normalize_workspace(raw: Any) -> tuple[Workspace, list[str]]:
    payload, format_name = unwrap_payload(raw)
    diagnostics = [f"识别格式：{format_name}"]
    projects_raw = payload.get("projects") or []
    if not projects_raw:
        raise ValueError("JSON 中没有项目。")
    projects: list[Project] = []
    for index, item in enumerate(projects_raw, start=1):
        tasks = [_normalize_task(task, payload) for task in item.get("tasks", [])]
        task_ids = {task.id for task in tasks}
        logs: list[DailyLog] = []
        for log_item in pick(item, "dailyLogs", "daily_logs", default=[]):
            if not isinstance(log_item, dict):
                continue
            log = _normalize_log(log_item)
            if log.result == "延期" and not log.delayReason:
                diagnostics.append(f"项目 {index} 有延期日报缺少原因，已标记为待补充。")
            if not task_ids or log.taskId in task_ids:
                logs.append(log)
        archives = [
            _normalize_archive(archive)
            for archive in pick(item, "archives", "archiveItems", "archive_items", default=[])
            if isinstance(archive, dict)
        ]
        project = Project(
            id=str(pick(item, "id", default=new_id())),
            name=str(pick(item, "name", "title", default="未命名项目")),
            deadline=clean_date(pick(item, "deadline")),
            summary=str(pick(item, "summary", default="")),
            topRisk=str(pick(item, "topRisk", "top_risk", default="")),
            nextStep=str(pick(item, "nextStep", "next_step", default="")),
            isPublic=bool(pick(item, "isPublic", "is_public", default=False)),
            publicSlug=str(pick(item, "publicSlug", "public_slug", default="")),
            tasks=tasks,
            dailyLogs=logs,
            archives=archives,
        )
        projects.append(project)
    selected_project_id = str(pick(payload, "selectedProjectId", "selected_project_id", default=projects[0].id))
    if selected_project_id not in {project.id for project in projects}:
        selected_project_id = projects[0].id
    inbox_tasks = [
        _normalize_inbox(item)
        for item in pick(payload, "inboxTasks", "inbox_tasks", "temporaryTasks", "temporary_tasks", default=[])
        if isinstance(item, dict)
    ]
    workspace = Workspace(
        selectedProjectId=selected_project_id,
        selectedDate=clean_date(pick(payload, "selectedDate", "selected_date")),
        projects=projects,
        inboxTasks=inbox_tasks,
        version=str(pick(payload, "version", default=APP_VERSION)),
    )
    diagnostics.append(f"归一化项目 {len(projects)} 个、待归档任务 {len(inbox_tasks)} 条。")
    return workspace, diagnostics


def _normalize_task(item: dict[str, Any], payload: dict[str, Any]) -> Task:
    risk = str(pick(item, "risk", default="M"))
    status = str(pick(item, "status", default="Open"))
    entries_raw = pick(item, "progressEntries", "progress_entries", default=[])
    entries = [
        ProgressEntry(
            entryDate=clean_date(pick(entry, "entryDate", "entry_date"), clean_date(pick(payload, "selectedDate", "selected_date"))),
            plannedProgress=clamp_progress(pick(entry, "plannedProgress", "planned_progress")),
            actualProgress=clamp_progress(pick(entry, "actualProgress", "actual_progress")),
        )
        for entry in entries_raw
        if isinstance(entry, dict)
    ]
    if not entries and ("plannedProgress" in item or "actualProgress" in item):
        entries.append(
            ProgressEntry(
                entryDate=clean_date(pick(payload, "selectedDate", "selected_date")),
                plannedProgress=clamp_progress(item.get("plannedProgress")),
                actualProgress=clamp_progress(item.get("actualProgress")),
            )
        )
    return Task(
        id=str(pick(item, "id", default=new_id())),
        parentId=pick(item, "parentId", "parent_id", default=None),
        risk=risk if risk in RISKS else "M",
        title=str(pick(item, "title", "task", "name", default="未命名任务")),
        responsible=str(pick(item, "responsible", "owner", default="")),
        startDate=clean_date(pick(item, "startDate", "start_date")),
        duration=max(1, parse_int(pick(item, "duration", default=1), 1)),
        status=status if status in STATUSES else "Open",
        completedDate=clean_date(pick(item, "completedDate", "completed_date"), ""),
        note=str(pick(item, "note", "remark", default="")),
        archivePath=str(pick(item, "archivePath", "archive_path", "filePath", "file_path", default="")),
        archiveType=str(pick(item, "archiveType", "archive_type", default="实验数据")),
        archiveKeywords=str(pick(item, "archiveKeywords", "archive_keywords", "keywords", default="")),
        progressEntries=entries,
    )


def _normalize_log(item: dict[str, Any]) -> DailyLog:
    return DailyLog(
        id=str(pick(item, "id", default=new_id())),
        date=clean_date(pick(item, "date", "log_date", "entry_date")),
        responsible=str(pick(item, "responsible", "owner", default="")),
        taskId=str(pick(item, "taskId", "task_id", default="")),
        planText=str(pick(item, "planText", "plan_text", default="")),
        actualText=str(pick(item, "actualText", "actual_text", default="")),
        plannedProgress=clamp_progress(pick(item, "plannedProgress", "planned_progress")),
        actualProgress=clamp_progress(pick(item, "actualProgress", "actual_progress", "progressAfter")),
        result=str(pick(item, "result", default="部分完成")),
        delayReason=str(pick(item, "delayReason", "delay_reason", default="")),
    )


def _normalize_archive(item: dict[str, Any]) -> ArchiveItem:
    item_type = str(pick(item, "type", "archiveType", "archive_type", default="其他"))
    return ArchiveItem(
        id=str(pick(item, "id", default=new_id())),
        date=clean_date(pick(item, "date", "archiveDate", "archive_date")),
        type=item_type if item_type in ARCHIVE_TYPES else "其他",
        title=str(pick(item, "title", "name", default="未命名档案")),
        owner=str(pick(item, "owner", "responsible", default="")),
        keywords=str(pick(item, "keywords", "tags", default="")),
        summary=str(pick(item, "summary", "note", default="")),
        path=str(pick(item, "path", "filePath", "file_path", default="")),
        relatedTaskId=str(pick(item, "relatedTaskId", "related_task_id", default="")),
        status=str(pick(item, "status", default="已归档")),
    )


def _normalize_inbox(item: dict[str, Any]) -> InboxTask:
    status = str(pick(item, "status", default="待归档"))
    return InboxTask(
        id=str(pick(item, "id", default=new_id())),
        createdDate=clean_date(pick(item, "createdDate", "created_date", "date")),
        title=str(pick(item, "title", "name", default="")),
        description=str(pick(item, "description", "note", "summary", default="")),
        source=str(pick(item, "source", default="")),
        status=status if status in INBOX_STATUSES else "待归档",
        suggestedAction=str(pick(item, "suggestedAction", "suggested_action", default="")),
        suggestedProjectId=str(pick(item, "suggestedProjectId", "suggested_project_id", default="")),
        suggestionReason=str(pick(item, "suggestionReason", "suggestion_reason", default="")),
        confirmed=bool(pick(item, "confirmed", default=False)),
    )


def load_workspace_json(path: Path) -> tuple[Workspace, list[str]]:
    return normalize_workspace(json.loads(path.read_text(encoding="utf-8-sig")))


def dump_workspace_json(workspace: Workspace, path: Path) -> None:
    payload = {"version": APP_VERSION, "workspace": to_dict(workspace)}
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def export_tasks_csv(project: Project, path: Path) -> None:
    rows = [["风险", "任务", "负责人", "开始", "工期", "结束", "状态", "计划%", "实际%", "实际完成日", "备注"]]
    for task in project.tasks:
        planned, actual = _latest(task)
        rows.append([task.risk, task.title, task.responsible, task.startDate, task.duration, task_end_date(task), task.status, planned, actual, task.completedDate, task.note])
    with path.open("w", encoding="utf-8-sig", newline="") as handle:
        csv.writer(handle).writerows(rows)


def export_project_excel(project: Project, path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "项目管理表"
    header_fill = PatternFill("solid", fgColor="E5E7EB")
    blue_fill = PatternFill("solid", fgColor="BFDBFE")
    green_fill = PatternFill("solid", fgColor="86EFAC")
    thin = Side(style="thin", color="9CA3AF")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    def style_row(row: int) -> None:
        for cell in ws[row]:
            cell.border = border
            cell.alignment = Alignment(vertical="top", wrap_text=True)

    planned, actual = project_progress(project)
    ws.append([project.name, "", "", ""])
    ws["A1"].font = Font(size=18, bold=True)
    ws.append(["Deadline", project.deadline, "计划进度", f"{planned}%"])
    ws.append(["一句话总结", project.summary, "实际进度", f"{actual}%"])
    ws.append(["TOP 风险", project.topRisk, "下一步计划", project.nextStep])
    for row in range(2, 5):
        style_row(row)

    ws.append([])
    task_header_row = ws.max_row + 1
    ws.append(["风险", "任务", "负责人", "开始", "工期", "结束", "状态", "计划%", "实际%", "实际完成日", "备注"])
    for cell in ws[task_header_row]:
        cell.fill = header_fill
        cell.font = Font(bold=True)
    style_row(task_header_row)
    for task in project.tasks:
        p, a = _latest(task)
        ws.append([task.risk, task.title, task.responsible, task.startDate, task.duration, task_end_date(task), task.status, p, a, task.completedDate, task.note])
        style_row(ws.max_row)

    ws.append([])
    log_header_row = ws.max_row + 1
    ws.append(["日期", "负责人", "关联任务", "计划完成", "实际完成", "计划%", "实际%", "结果", "延期原因"])
    for cell in ws[log_header_row]:
        cell.fill = header_fill
        cell.font = Font(bold=True)
    style_row(log_header_row)
    task_names = {task.id: task.title for task in project.tasks}
    for log in project.dailyLogs:
        ws.append([log.date, log.responsible, task_names.get(log.taskId, ""), log.planText, log.actualText, log.plannedProgress, log.actualProgress, log.result, log.delayReason])
        style_row(ws.max_row)

    ws.append([])
    archive_header_row = ws.max_row + 1
    ws.append(["档案日期", "类型", "标题", "负责人", "关键词", "摘要", "路径", "状态"])
    for cell in ws[archive_header_row]:
        cell.fill = header_fill
        cell.font = Font(bold=True)
    style_row(archive_header_row)
    for archive in project.archives:
        ws.append([archive.date, archive.type, archive.title, archive.owner, archive.keywords, archive.summary, archive.path, archive.status])
        style_row(ws.max_row)

    ws.append([])
    dates = _gantt_dates(project)
    gantt_header_row = ws.max_row + 1
    ws.append(["任务", "负责人", "风险", "实际%", *[item[5:] for item in dates]])
    for cell in ws[gantt_header_row]:
        cell.fill = header_fill
        cell.font = Font(bold=True)
    style_row(gantt_header_row)
    for task in project.tasks:
        _, actual_progress = _latest(task)
        row = [task.title, task.responsible, task.risk, f"{actual_progress}%"]
        end = task_end_date(task)
        row.extend("■" if task.startDate <= item <= end else "" for item in dates)
        ws.append(row)
        for index, date_value in enumerate(dates, start=5):
            if task.startDate <= date_value <= end:
                ws.cell(ws.max_row, index).fill = green_fill if task.completedDate and date_value <= task.completedDate else blue_fill
        style_row(ws.max_row)

    widths = {1: 22, 2: 34, 3: 16, 4: 14, 5: 16, 6: 36, 7: 44, 8: 14, 9: 14, 10: 14, 11: 32}
    for index, width in widths.items():
        ws.column_dimensions[get_column_letter(index)].width = width
    for index in range(12, ws.max_column + 1):
        ws.column_dimensions[get_column_letter(index)].width = 5
    wb.save(path)


def week_bounds(selected_date: str | None) -> tuple[str, str]:
    normalized = normalize_date(selected_date) or today()
    current = date.fromisoformat(normalized)
    start = current - timedelta(days=current.weekday())
    end = start + timedelta(days=6)
    return start.isoformat(), end.isoformat()


def export_weekly_report_excel(project: Project, week_start: str, week_end: str, path: Path) -> None:
    wb = Workbook()
    header_fill = PatternFill("solid", fgColor="E5E7EB")
    section_fill = PatternFill("solid", fgColor="DBEAFE")
    warning_fill = PatternFill("solid", fgColor="FEE2E2")
    thin = Side(style="thin", color="CBD5E1")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    task_names = {task.id: task.title for task in project.tasks}
    planned, actual = project_progress(project)

    def setup_sheet(ws, widths: list[int]) -> None:
        ws.freeze_panes = "A2"
        for index, width in enumerate(widths, start=1):
            ws.column_dimensions[get_column_letter(index)].width = width

    def style_row(ws, row: int, fill=None, bold: bool = False) -> None:
        for cell in ws[row]:
            cell.border = border
            cell.alignment = Alignment(vertical="top", wrap_text=True)
            if fill:
                cell.fill = fill
            if bold:
                cell.font = Font(bold=True)

    def append_header(ws, values: list[str]) -> None:
        ws.append(values)
        style_row(ws, ws.max_row, header_fill, True)

    ws = wb.active
    ws.title = "周报摘要"
    setup_sheet(ws, [18, 48, 18, 48])
    ws.append([project.name, "", "周报周期", f"{week_start} 至 {week_end}"])
    ws["A1"].font = Font(size=18, bold=True)
    style_row(ws, 1, section_fill, True)
    deadline_delta = _deadline_delta_text(project.deadline, week_end)
    summary_rows = [
        ["Deadline", project.deadline, "Deadline 状态", deadline_delta],
        ["计划进度", f"{planned}%", "实际进度", f"{actual}%"],
        ["一句话总结", project.summary, "", ""],
        ["TOP 风险", project.topRisk, "", ""],
        ["下一步计划", project.nextStep, "", ""],
    ]
    for row in summary_rows:
        ws.append(row)
        style_row(ws, ws.max_row)

    logs = _weekly_logs(project, week_start, week_end)
    ws = wb.create_sheet("本周日报")
    setup_sheet(ws, [14, 16, 26, 36, 36, 12, 12, 14, 34])
    append_header(ws, ["日期", "负责人", "关联任务", "计划完成", "实际完成", "计划%", "实际%", "结果", "延期原因"])
    for log in logs:
        ws.append([log.date, log.responsible, task_names.get(log.taskId, ""), log.planText, log.actualText, log.plannedProgress, log.actualProgress, log.result, log.delayReason])
        style_row(ws, ws.max_row, warning_fill if log.result == "延期" else None)
    if not logs:
        ws.append(["本周暂无日报记录"])
        style_row(ws, ws.max_row)

    ws = wb.create_sheet("任务进展")
    setup_sheet(ws, [10, 34, 14, 14, 14, 10, 14, 12, 12, 34])
    append_header(ws, ["风险", "任务", "负责人", "状态", "开始", "工期", "结束", "计划%", "实际%", "备注"])
    for task in _sorted_tasks(project.tasks, week_end):
        p, a = latest_progress(task, week_end)
        overdue = _is_task_overdue(task, week_end)
        ws.append([task.risk, task.title, task.responsible, _status_label(task.status), task.startDate, task.duration, task_end_date(task), p, a, task.note])
        style_row(ws, ws.max_row, warning_fill if overdue else None)

    ws = wb.create_sheet("风险与延期")
    setup_sheet(ws, [14, 14, 30, 16, 48, 40])
    append_header(ws, ["类型", "日期", "任务", "负责人", "说明", "延期原因"])
    risk_rows = _risk_and_delay_rows(project, week_start, week_end)
    for row in risk_rows:
        ws.append(row)
        style_row(ws, ws.max_row, warning_fill)
    if not risk_rows:
        ws.append(["本周无高风险、延期日报或逾期未关闭任务"])
        style_row(ws, ws.max_row)

    ws = wb.create_sheet("下周计划")
    setup_sheet(ws, [12, 34, 16, 14, 14, 14, 42])
    ws.append(["项目下一步计划", project.nextStep])
    style_row(ws, ws.max_row, section_fill, True)
    ws.append([])
    append_header(ws, ["风险", "任务", "负责人", "状态", "开始", "结束", "备注"])
    for task in _next_week_tasks(project.tasks, week_end):
        ws.append([task.risk, task.title, task.responsible, _status_label(task.status), task.startDate, task_end_date(task), task.note])
        style_row(ws, ws.max_row)

    ws = wb.create_sheet("归档证据")
    setup_sheet(ws, [14, 14, 30, 16, 24, 44, 52, 14])
    append_header(ws, ["档案日期", "类型", "标题", "负责人", "关键词", "摘要", "路径", "状态"])
    for archive in sorted(project.archives, key=lambda item: item.date):
        ws.append([archive.date, archive.type, archive.title, archive.owner, archive.keywords, archive.summary, archive.path, archive.status])
        style_row(ws, ws.max_row)
    if not project.archives:
        ws.append(["当前项目暂无归档证据"])
        style_row(ws, ws.max_row)

    wb.save(path)


def export_project_briefing_ppt(project: Project, week_start: str, week_end: str, path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    planned, actual = project_progress(project)
    overdue = sum(1 for task in project.tasks if _is_task_overdue(task, week_end))
    logs = _weekly_logs(project, week_start, week_end)

    def add_slide(title: str, subtitle: str = ""):
        slide = prs.slides.add_slide(blank)
        _add_text(slide, title, 0.55, 0.38, 8.8, 0.45, 24, bold=True, color=RGBColor(15, 23, 42))
        if subtitle:
            _add_text(slide, subtitle, 0.58, 0.88, 8.4, 0.28, 10, color=RGBColor(100, 116, 139))
        _add_text(slide, f"{project.name} · {week_start} 至 {week_end}", 9.4, 0.45, 3.2, 0.28, 10, align=PP_ALIGN.RIGHT, color=RGBColor(100, 116, 139))
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.2), Inches(12.25), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(226, 232, 240)
        line.line.fill.background()
        return slide

    slide = add_slide("项目周报概览", project.summary or "本周项目状态汇总")
    _add_metric(slide, "Deadline", project.deadline, _deadline_delta_text(project.deadline, week_end), 0.65, 1.55)
    _add_metric(slide, "计划进度", f"{planned}%", "基于任务最新计划进度", 3.75, 1.55)
    _add_metric(slide, "实际进度", f"{actual}%", "基于任务最新实际进度", 6.85, 1.55, RGBColor(15, 118, 110))
    _add_metric(slide, "逾期任务", str(overdue), "截至本周末未关闭", 9.95, 1.55, RGBColor(220, 38, 38) if overdue else RGBColor(15, 118, 110))
    _add_panel(slide, "一句话总结", project.summary or "暂无一句话总结。", 0.65, 3.05, 3.8, 2.15, RGBColor(37, 99, 235))
    _add_panel(slide, "TOP 风险", project.topRisk or "暂无 TOP 风险。", 4.75, 3.05, 3.8, 2.15, RGBColor(220, 38, 38))
    _add_panel(slide, "下一步计划", project.nextStep or "暂无下一步计划。", 8.85, 3.05, 3.8, 2.15, RGBColor(15, 118, 110))

    slide = add_slide("本周完成与日报摘要")
    completed = [log for log in logs if log.result == "完成" or int(log.actualProgress) >= 100]
    partial = [log for log in logs if log not in completed and log.result != "延期"]
    _add_table(
        slide,
        ["日期", "任务", "负责人", "实际完成", "实际%"],
        [[log.date, _task_title(project, log.taskId), log.responsible, log.actualText, f"{log.actualProgress}%"] for log in completed[:8]] or [["-", "本周暂无完成记录", "", "", ""]],
        0.65,
        1.55,
        5.95,
        4.8,
    )
    _add_table(
        slide,
        ["日期", "任务", "负责人", "计划/实际", "结果"],
        [[log.date, _task_title(project, log.taskId), log.responsible, f"{log.plannedProgress}% / {log.actualProgress}%", log.result] for log in partial[:8]] or [["-", "本周暂无部分完成记录", "", "", ""]],
        6.9,
        1.55,
        5.75,
        4.8,
    )

    slide = add_slide("任务进展")
    task_rows = []
    for task in _sorted_tasks(project.tasks, week_end)[:10]:
        p, a = latest_progress(task, week_end)
        task_rows.append([task.risk, task.title, task.responsible, _status_label(task.status), task_end_date(task), f"{p}% / {a}%"])
    _add_table(slide, ["风险", "任务", "负责人", "状态", "结束", "计划/实际"], task_rows or [["-", "暂无任务", "", "", "", ""]], 0.65, 1.55, 12.0, 4.95)

    slide = add_slide("风险与延期")
    risk_rows = _risk_and_delay_rows(project, week_start, week_end)
    _add_table(slide, ["类型", "日期", "任务", "负责人", "说明", "原因"], risk_rows[:9] or [["-", "-", "暂无高风险、延期或逾期项", "", "", ""]], 0.65, 1.55, 12.0, 4.95)

    slide = add_slide("下周计划")
    _add_panel(slide, "项目下一步计划", project.nextStep or "暂无下一步计划。", 0.65, 1.45, 12.0, 1.4, RGBColor(15, 118, 110))
    next_rows = [[task.risk, task.title, task.responsible, _status_label(task.status), task_end_date(task), task.note] for task in _next_week_tasks(project.tasks, week_end)[:8]]
    _add_table(slide, ["风险", "任务", "负责人", "状态", "结束", "备注"], next_rows or [["-", "暂无未关闭任务", "", "", "", ""]], 0.65, 3.15, 12.0, 3.0)

    prs.save(path)


def _latest(task: Task) -> tuple[int, int]:
    if not task.progressEntries:
        return 0, 0
    entry = sorted(task.progressEntries, key=lambda item: item.entryDate)[-1]
    return entry.plannedProgress, entry.actualProgress


def _weekly_logs(project: Project, week_start: str, week_end: str) -> list[DailyLog]:
    return sorted(
        [log for log in project.dailyLogs if week_start <= (normalize_date(log.date) or log.date) <= week_end],
        key=lambda item: ((normalize_date(item.date) or item.date), item.responsible, item.taskId),
    )


def _status_label(status: str) -> str:
    return {"Open": "未开始", "Ongoing": "进行中", "Closed": "已关闭"}.get(status, status)


def _risk_rank(risk: str) -> int:
    return {"H": 0, "M": 1, "L": 2}.get(risk, 3)


def _sorted_tasks(tasks: list[Task], on_date: str) -> list[Task]:
    return sorted(tasks, key=lambda task: (task.status == "Closed", _risk_rank(task.risk), task_end_date(task), task.title))


def _is_task_overdue(task: Task, on_date: str) -> bool:
    return task.status != "Closed" and task_end_date(task) < on_date


def _deadline_delta_text(deadline: str, on_date: str) -> str:
    normalized = normalize_date(deadline) or deadline
    try:
        delta = (date.fromisoformat(normalized) - date.fromisoformat(on_date)).days
    except ValueError:
        return "Deadline 无效"
    if delta >= 0:
        return f"剩余 {delta} 天"
    return f"已逾期 {abs(delta)} 天"


def _task_title(project: Project, task_id: str) -> str:
    task = next((item for item in project.tasks if item.id == task_id), None)
    return task.title if task else ""


def _risk_and_delay_rows(project: Project, week_start: str, week_end: str) -> list[list[str]]:
    rows: list[list[str]] = []
    task_by_id = {task.id: task for task in project.tasks}
    for task in _sorted_tasks(project.tasks, week_end):
        if task.risk == "H":
            rows.append(["高风险任务", "", task.title, task.responsible, f"{_status_label(task.status)}，结束 {task_end_date(task)}", task.note])
        if _is_task_overdue(task, week_end):
            rows.append(["逾期未关闭", task_end_date(task), task.title, task.responsible, f"截至 {week_end} 未关闭", task.note])
    for log in _weekly_logs(project, week_start, week_end):
        if log.result == "延期":
            task = task_by_id.get(log.taskId)
            rows.append(["延期日报", log.date, task.title if task else "", log.responsible, log.actualText or log.planText, log.delayReason])
    return rows


def _next_week_tasks(tasks: list[Task], week_end: str) -> list[Task]:
    next_start = add_days(week_end, 1)
    next_end = add_days(week_end, 7)
    candidates = [
        task
        for task in tasks
        if task.status != "Closed" and (task.startDate <= next_end or task_end_date(task) >= next_start)
    ]
    return _sorted_tasks(candidates, next_end)


def _add_text(slide, text: str, left: float, top: float, width: float, height: float, size: int = 12, bold: bool = False, color: RGBColor | None = None, align=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = _clip(text, 240)
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or RGBColor(30, 41, 59)
    return box


def _add_metric(slide, label: str, value: str, note: str, left: float, top: float, color: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.75), Inches(1.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
    shape.line.color.rgb = RGBColor(226, 232, 240)
    _add_text(slide, label, left + 0.18, top + 0.12, 2.25, 0.22, 9, bold=True, color=RGBColor(100, 116, 139))
    _add_text(slide, value, left + 0.18, top + 0.38, 2.25, 0.32, 20, bold=True, color=color or RGBColor(15, 23, 42))
    _add_text(slide, note, left + 0.18, top + 0.78, 2.25, 0.18, 8, color=RGBColor(100, 116, 139))


def _add_panel(slide, title: str, body: str, left: float, top: float, width: float, height: float, accent: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(226, 232, 240)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.05), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    _add_text(slide, title, left + 0.22, top + 0.2, width - 0.35, 0.28, 12, bold=True, color=RGBColor(15, 23, 42))
    _add_text(slide, body, left + 0.22, top + 0.65, width - 0.35, height - 0.82, 12, color=RGBColor(30, 41, 59))


def _add_table(slide, headers: list[str], rows: list[list[str]], left: float, top: float, width: float, height: float):
    row_count = max(2, len(rows) + 1)
    col_count = len(headers)
    table_shape = slide.shapes.add_table(row_count, col_count, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(241, 245, 249)
        _format_cell(cell, bold=True)
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row[:col_count]):
            cell = table.cell(row_index, col_index)
            cell.text = _clip(str(value or ""), 70)
            _format_cell(cell)


def _format_cell(cell, bold: bool = False):
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(8 if not bold else 8.5)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(30, 41, 59)
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)


def _clip(text: str, limit: int) -> str:
    text = " ".join(str(text or "").split())
    return text if len(text) <= limit else text[: limit - 1] + "…"


def _gantt_dates(project: Project) -> list[str]:
    if not project.tasks:
        return [add_days(today(), index) for index in range(14)]
    start = min(task.startDate for task in project.tasks)
    end = max(task_end_date(task) for task in project.tasks)
    total = min(max((__import__("datetime").date.fromisoformat(end) - __import__("datetime").date.fromisoformat(start)).days + 1, 14), 120)
    return [add_days(start, index) for index in range(total)]
